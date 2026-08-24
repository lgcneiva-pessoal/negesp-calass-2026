# -*- coding: utf-8 -*-
"""Renderiza os 22 slides do deck HTML e monta PDF + PPTX fiéis ao que está no ar.

Uso (com o repo servido em http://localhost:8791):
    python3 source/render_deck.py "http://localhost:8791/v2/index.html"    ~/Downloads/Negesp-CALASS-2026-PT
    python3 source/render_deck.py "http://localhost:8791/v2-fr/index.html" ~/Downloads/Negesp-CALASS-2026-FR

Requer: playwright (+chromium), Pillow, python-pptx.
Ver HANDOFF.md secao 11 para os detalhes.
"""
import sys, os, io
from playwright.sync_api import sync_playwright
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

URL, OUTBASE = sys.argv[1], sys.argv[2]
N, W, H, SCALE = 22, 1280, 720, 2
TMP = os.path.abspath(os.path.join(os.path.dirname(os.path.abspath(__file__)),
                                   '..', '.render', os.path.basename(OUTBASE)))
os.makedirs(TMP, exist_ok=True)

# CSS de exportacao: reaproveita as regras do @media print do proprio deck
# (forcam TODAS as revelacoes completas) sem quebrar o isolamento por slide.
EXPORT_CSS = """
/* revelacoes: estado final garantido (evita capturar no meio da animacao) */
/* NAO usar transform:none aqui: varios elementos (#s11 .reading, #s10 .fn...) dependem
   de translate(-50%,-50%) para se posicionar, e zera-lo quebra o layout. O proprio deck
   ja resolve o transform da revelacao quando o slide esta .on / o passo esta .go. */
.rv,.st{opacity:1!important;filter:none!important}
#s01 h1 .w,#s02 .q .w,#s19 .msg .w{opacity:1!important;transform:none!important;filter:none!important}
/* a regra acima reacende o #orb5 (que e .rv) por cima dos cards de resumo.
   visibility (nao display:none): esconde mas MANTEM o espaco no fluxo, senao o
   sobretitulo do s13 escorrega para o meio e some atras dos cards. */
#s13.summaryon #orb5{visibility:hidden!important}
/* mapas: preenchimento de impressao (o de tela fica palido demais no papel) */
.brmap path.uf{fill:rgba(14,122,78,.25)!important;stroke:#2FBF83!important}
.brmap path.uf.dest{fill:rgba(255,138,60,.65)!important}
/* arcos da rede: no vivo pulsam em ondas; no estatico ficam todos desenhados */
#s16 .brmap path.net{opacity:1!important;stroke-dashoffset:0!important;stroke:rgba(14,122,78,.55)!important}
/* cromo de tela que nao faz sentido impresso */
.notesbar,.menu{display:none!important}
.hud .hbl{display:none!important}
#s02 .caret{display:none!important}          /* cursor piscante congelado */
#s10 .bpulse{display:none!important}         /* pulso que corre o conector */
#s11 .sweep{display:none!important}          /* varredura diagonal giratoria */
"""

# congela animacoes: finitas vao ao estado final, infinitas voltam ao inicio limpo
FREEZE_JS = """() => {
  document.getAnimations().forEach(a => {
    try {
      const t = (a.effect && a.effect.getTiming) ? a.effect.getTiming() : {};
      if (t.iterations === Infinity || t.duration === Infinity) { a.currentTime = 0; a.pause(); }
      else { a.finish(); }
    } catch (e) {}
  });
}"""

frames = []
with sync_playwright() as p:
    br = p.chromium.launch(args=['--force-color-profile=srgb', '--font-render-hinting=none'])
    pg = br.new_page(viewport={'width': W, 'height': H}, device_scale_factor=SCALE)
    pg.goto(URL, wait_until='networkidle')
    pg.add_style_tag(content=EXPORT_CSS)
    pg.wait_for_timeout(1500)                      # fonte embutida + primeiro layout

    for i in range(N):
        pg.evaluate("i => window.goTo(i, true)", i)
        pg.wait_for_timeout(400)
        sid = pg.evaluate("() => (document.querySelector('.slide.on')||{}).id")

        if sid in ('s11', 's16'):                  # mapas: o JS precisa CRIAR os paths antes
            pg.wait_for_timeout(2600)
        elif sid == 's13':                         # orbita: mostra os 4 cards de resumo
            pg.evaluate("() => window.__s13step && window.__s13step(5)")
            pg.wait_for_timeout(1200)
        else:
            pg.wait_for_timeout(900)

        pg.evaluate(FREEZE_JS)                     # so entao congela tudo
        pg.wait_for_timeout(350)

        f = os.path.join(TMP, f'slide-{i+1:02d}.png')
        pg.screenshot(path=f)
        frames.append(f)
        print(f'  slide {i+1:02d}/{N} ({sid})', flush=True)
    br.close()

imgs = [Image.open(f).convert('RGB') for f in frames]
imgs[0].save(OUTBASE + '.pdf', save_all=True, append_images=imgs[1:], resolution=(W*SCALE)/13.3333)
print('PDF :', OUTBASE + '.pdf', f'{os.path.getsize(OUTBASE + ".pdf")//1024} KB')

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.3333), Inches(7.5)
for im in imgs:
    s = prs.slides.add_slide(prs.slide_layouts[6])
    buf = io.BytesIO(); im.save(buf, 'JPEG', quality=93, optimize=True); buf.seek(0)
    s.shapes.add_picture(buf, 0, 0, width=prs.slide_width, height=prs.slide_height)
prs.save(OUTBASE + '.pptx')
print('PPTX:', OUTBASE + '.pptx', f'{os.path.getsize(OUTBASE + ".pptx")//1024} KB')
